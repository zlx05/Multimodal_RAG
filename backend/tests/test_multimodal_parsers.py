from pathlib import Path

from backend.app.rag.blocks import DocumentBlock
from backend.app.rag.parsers import create_parser, get_parser_type
from backend.app.rag.parsers.docx_parser import DocxParser
from backend.app.rag.parsers.image_parser import ImageParser
from backend.app.rag.parsers.html_parser import HtmlParser
from backend.app.rag.parsers.pptx_parser import PptxParser
from backend.app.rag.vision import VisionAnalysis, VisionAnalyzer
from backend.app.rag.ocr.base import OcrResult
from backend.app.rag.parsers.media import _is_plausible_formula


class FakeOcr:
    def __call__(self, path):
        return OcrResult(
            text="OCR 识别出的公式 F = ma",
            boxes=[[0, 0, 10, 0, 10, 10, 0, 10]],
            confidences=[0.91],
            confidence=0.91,
        )


class FakeVision:
    def analyze(self, path, ocr_text=""):
        return VisionAnalysis(
            text="这是一道关于牛顿第二定律的手写题。",
            metadata={"vision_model": "fake-vision"},
        )


def test_image_parser_prefers_vision_and_keeps_ocr_evidence(tmp_path):
    image = tmp_path / "note.png"
    image.write_bytes(b"fake image")

    blocks = ImageParser(
        "doc-1",
        ocr_engine=FakeOcr(),
        vision_analyzer=FakeVision(),
        original_dir=str(tmp_path / "original"),
    ).parse(image)

    assert len(blocks) == 1
    assert blocks[0].content_type == "image_description"
    assert "牛顿第二定律" in blocks[0].text
    assert blocks[0].metadata["ocr_boxes"]
    assert blocks[0].metadata["ocr_text"] == "OCR 识别出的公式 F = ma"
    assert blocks[0].metadata["vision_model"] == "fake-vision"


def test_docx_parser_extracts_heading_and_table(tmp_path):
    from docx import Document

    path = tmp_path / "lesson.docx"
    document = Document()
    document.add_heading("力学", level=1)
    document.add_paragraph("牛顿第二定律")
    # Native Word math (OMML) should be preserved as LaTeX, not OCR text.
    from docx.oxml import OxmlElement

    formula_paragraph = document.add_paragraph()
    math = OxmlElement("m:oMath")
    run = OxmlElement("m:r")
    math_text = OxmlElement("m:t")
    math_text.text = "x"
    run.append(math_text)
    math.append(run)
    formula_paragraph._p.append(math)
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "公式"
    table.cell(0, 1).text = "F=ma"
    table.cell(1, 0).text = "单位"
    table.cell(1, 1).text = "N"
    document.save(path)

    blocks = DocxParser("doc-2", ocr_engine=FakeOcr(), work_dir=str(tmp_path)).parse(path)

    assert any(block.content_type == "heading" for block in blocks)
    table_block = next(block for block in blocks if block.content_type == "table")
    assert "F=ma" in table_block.text
    assert table_block.metadata["table"][1][1] == "N"
    formula_block = next(block for block in blocks if block.content_type == "formula")
    assert "$$" in formula_block.text
    assert formula_block.metadata["formulas_latex"] == ["x"]


def test_pptx_parser_extracts_slide_text_and_table(tmp_path):
    from pptx import Presentation

    path = tmp_path / "lesson.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, 3000000, 500000)
    textbox.text = "动量守恒"
    table = slide.shapes.add_table(2, 2, 0, 600000, 3000000, 800000).table
    table.cell(0, 0).text = "定律"
    table.cell(0, 1).text = "动量守恒"
    presentation.save(path)

    blocks = PptxParser("doc-3", ocr_engine=FakeOcr(), work_dir=str(tmp_path)).parse(path)

    assert any("动量守恒" in block.text and block.page_number == 1 for block in blocks)
    assert any(block.content_type == "table" for block in blocks)


def test_parser_factory_supports_office_formats():
    assert get_parser_type("legacy.doc") == "legacy_doc"
    assert get_parser_type("legacy.ppt") == "legacy_ppt"
    assert get_parser_type("lesson.docx") == "docx"
    assert get_parser_type("lesson.pptx") == "pptx"
    assert isinstance(create_parser("doc-4", "lesson.docx", ocr_engine=FakeOcr()), DocxParser)
    assert isinstance(create_parser("doc-5", "lesson.pptx", ocr_engine=FakeOcr()), PptxParser)


def test_office_parsers_accept_worker_kwargs(tmp_path):
    """worker 对所有解析器统一传 original_dir/work_dir/vision 等 kwargs；
    docx/pptx 必须像 pdf/html/image 一样接收 original_dir，否则上传直接崩
    （回归：PptxParser.__init__ got an unexpected keyword argument 'original_dir'）。"""
    kwargs = {
        "original_dir": str(tmp_path / "original"),
        "work_dir": str(tmp_path / "work"),
        "ocr_engine": FakeOcr(),
    }
    assert isinstance(create_parser("doc-6", "lesson.docx", **kwargs), DocxParser)
    assert isinstance(create_parser("doc-7", "lesson.pptx", **kwargs), PptxParser)


def test_vision_analyzer_uses_openai_compatible_image_message(tmp_path):
    class FakeCompletions:
        def __init__(self):
            self.kwargs = None

        def create(self, **kwargs):
            self.kwargs = kwargs
            return type(
                "Response",
                (),
                {"choices": [type("Choice", (), {"message": type("Message", (), {"content": "识别结果"})()})()]},
            )()

    class FakeClient:
        def __init__(self):
            self.chat = type("Chat", (), {"completions": FakeCompletions()})()

    image = tmp_path / "image.png"
    image.write_bytes(b"image")
    client = FakeClient()
    result = VisionAnalyzer(
        api_key="test-key",
        base_url="https://example.test/v1",
        model="vision-test",
        client=client,
    ).analyze(image)

    assert result is not None
    assert result.text == "识别结果"
    content = client.chat.completions.kwargs["messages"][0]["content"]
    assert content[1]["type"] == "image_url"
    assert content[1]["image_url"]["url"].startswith("data:image/png;base64,")


def test_formula_quality_gate_rejects_full_page_repetition():
    assert _is_plausible_formula(r"\frac{\sin x-x}{x^3}")
    assert not _is_plausible_formula(r"\begin{aligned}" + "\\quad 本题" * 20)


def test_html_parser_keeps_article_structure_and_image_alt(tmp_path):
    html = tmp_path / "article.html"
    html.write_text(
        "<html><body><nav>导航</nav><main><h1>Go 基本语法</h1>"
        "<p>变量和函数是基础。</p><img src=\"/missing.png\" alt=\"示例代码截图\" />"
        "</main></body></html>",
        encoding="utf-8",
    )
    blocks = HtmlParser(
        "doc-html",
        ocr_engine=FakeOcr(),
        work_dir=str(tmp_path),
        base_url="https://example.test/guide/",
    ).parse(html)

    assert any(block.content_type == "heading" and "Go 基本语法" in block.text for block in blocks)
    assert any(block.content_type == "image_description" and "示例代码截图" in block.text for block in blocks)


def test_html_parser_merges_list_items_into_one_block(tmp_path):
    """同一 ul/ol 的相邻 li 合并为一个文本块，避免"可比较类型有：布尔/数字/字符串"被切碎。"""
    html = tmp_path / "list.html"
    html.write_text(
        "<html><body><main><h2>比较</h2>"
        "<p>它们的参数支持所有的可比较类型，go 中的可比较类型有</p>"
        "<ul><li>布尔</li><li>数字</li><li>字符串</li><li>通道 （仅支持判断是否相等）</li></ul>"
        "<p>除此之外，还可以通过标准库 cmp 判断。</p>"
        "</main></body></html>",
        encoding="utf-8",
    )
    blocks = HtmlParser(
        "doc-list", ocr_engine=FakeOcr(), work_dir=str(tmp_path)
    ).parse(html)

    texts = [b.text for b in blocks]
    # 引导段 + 列表项并入同一个块，列表块自带检索词（可比较类型/比较），避免孤立不可召回
    lead_block = next(t for t in texts if "它们的参数支持所有的可比较类型" in t)
    assert "它们的参数支持所有的可比较类型，go 中的可比较类型有：布尔；数字；字符串；通道 （仅支持判断是否相等）" in lead_block
    # 列表后的段落仍是独立块，不被列表吞并
    assert any("标准库 cmp" in t for t in texts)
    assert not any(t == "它们的参数支持所有的可比较类型，go 中的可比较类型有" for t in texts)
