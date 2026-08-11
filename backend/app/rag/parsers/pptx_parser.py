"""PPTX parser for slide text, tables and embedded pictures."""

from __future__ import annotations

import json
import os
from pathlib import Path
from typing import Any, Iterable

from ..blocks import DocumentBlock
from ..ocr import create_ocr_engine
from ..vision import VisionAnalyzer
from .base import BaseParser
from .formula import append_formulas, extract_native_formulas
from .media import analyze_media


class PptxParser(BaseParser):
    source_type = "pptx"

    def __init__(
        self,
        document_id: str,
        ocr_engine=None,
        vision_analyzer: VisionAnalyzer | None = None,
        formula_recognizer=None,
        work_dir: str | None = None,
        original_dir: str | None = None,
    ):
        super().__init__(document_id)
        # original_dir 兼容统一接口（worker 对所有解析器传该参数）；PPTX 原文件在管理目录，无需复制
        self.ocr_engine = ocr_engine or create_ocr_engine()
        self.vision_analyzer = vision_analyzer
        self.formula_recognizer = formula_recognizer
        self.work_dir = Path(work_dir or os.getenv("RAG_WORK_DIR", "data/.work"))

    def parse(self, path: str | Path) -> list[DocumentBlock]:
        from pptx import Presentation

        presentation = Presentation(str(path))
        blocks: list[DocumentBlock] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            # 幻灯片标题作为 heading_path：分组时以幻灯片为单元（配合 slide_number 边界），
            # 标题也进入 search_text 的上下文，主题推断更准；无标题则保持空路径。
            heading: list[str] = []
            try:
                title_shape = slide.shapes.title
                title_text = title_shape.text.strip() if title_shape is not None else ""
            except Exception:
                title_text = ""
            if title_text:
                heading = [title_text]
            text_parts: list[str] = []
            for shape in self._iter_shapes(slide.shapes):
                if getattr(shape, "has_text_frame", False):
                    formulas = extract_native_formulas(shape._element)
                    text = append_formulas(shape.text, formulas)
                    if text:
                        text_parts.append(text)
                if getattr(shape, "has_table", False):
                    rows = []
                    formula_rows = []
                    for row in shape.table.rows:
                        row_values = []
                        row_formulas = []
                        for cell in row.cells:
                            formulas = extract_native_formulas(cell._tc)
                            row_values.append(append_formulas(cell.text, formulas))
                            row_formulas.append(formulas)
                        rows.append(row_values)
                        formula_rows.append(row_formulas)
                    if rows:
                        blocks.append(
                            self._block(
                                "PPT 表格：\n" + json.dumps(rows, ensure_ascii=False),
                                content_type="table",
                                page_number=slide_number,
                                heading_path=list(heading),
                                metadata={
                                    "slide_number": slide_number,
                                    "table": rows,
                                    "formulas_latex": formula_rows,
                                },
                            )
                        )
                if getattr(shape, "shape_type", None) == 13:  # MSO_SHAPE_TYPE.PICTURE
                    blocks.append(self._parse_picture(shape, slide_number, len(blocks), heading))

            if text_parts:
                blocks.append(
                    self._block(
                        "\n".join(text_parts),
                        page_number=slide_number,
                        heading_path=list(heading),
                        metadata={"slide_number": slide_number},
                    ),
                )
        return blocks

    def _parse_picture(self, shape: Any, slide_number: int, index: int, heading: list[str]) -> DocumentBlock:
        image_dir = self.work_dir / self.document_id / "pptx_images"
        image_path = image_dir / f"slide_{slide_number}_image_{index}.png"
        image_path.parent.mkdir(parents=True, exist_ok=True)
        image_path.write_bytes(shape.image.blob)
        text, metadata, confidence, content_type = analyze_media(
            image_path,
            self.ocr_engine,
            self.vision_analyzer,
            self.formula_recognizer,
        )
        metadata.update({"slide_number": slide_number, "asset_kind": "pptx_image"})
        return self._block(
            text,
            content_type=content_type,
            page_number=slide_number,
            heading_path=list(heading),
            image_path=str(image_path),
            confidence=confidence,
            metadata=metadata,
        )

    @staticmethod
    def _iter_shapes(shapes: Iterable[Any]):
        for shape in shapes:
            if getattr(shape, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP
                yield from PptxParser._iter_shapes(shape.shapes)
            else:
                yield shape
